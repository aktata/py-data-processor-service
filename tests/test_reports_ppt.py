from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation

from app.reporting.ppt_report import export_ppt_report


def test_export_ppt(tmp_path: Path) -> None:
    metrics = pd.DataFrame(
        [
            {
                "company_name": "Alpha",
                "year": 2023,
                "indicator_name": "net_profit_margin",
                "indicator_value": 0.2,
                "risk_level": "low",
                "risk_score": 10,
            },
            {
                "company_name": "Alpha",
                "year": 2023,
                "indicator_name": "current_ratio",
                "indicator_value": 1.5,
                "risk_level": "low",
                "risk_score": 15,
            },
            {
                "company_name": "Alpha",
                "year": 2023,
                "indicator_name": "roe",
                "indicator_value": 0.1,
                "risk_level": "medium",
                "risk_score": 55,
            },
        ]
    )
    overall = pd.DataFrame(
        [{"company_name": "Alpha", "year": 2023, "overall_risk_score": 20}]
    )
    ranking = metrics[metrics["indicator_name"] == "net_profit_margin"]

    output_path = tmp_path / "report.pptx"
    assets_dir = tmp_path / "assets"
    template = tmp_path / "template.pptx"
    Presentation().save(template)

    export_ppt_report(metrics, overall, output_path, assets_dir, template, ranking)
    assert output_path.exists()
    presentation = Presentation(output_path)
    assert len(presentation.slides) >= 2
