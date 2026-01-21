from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.reporting.charts import plot_bar, plot_trend

RISK_RGB = {
    "low": RGBColor(198, 239, 206),
    "medium": RGBColor(255, 235, 156),
    "high": RGBColor(255, 199, 206),
}


def _add_title(slide, text: str) -> None:
    title = slide.shapes.title
    if title is not None:
        title.text = text


def _add_textbox(slide, text: str, left: float, top: float, width: float, height: float) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(14)


def _add_indicator_table(slide, indicators: pd.DataFrame) -> None:
    rows = len(indicators) + 1
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(2.5)).table
    table.cell(0, 0).text = "指标"
    table.cell(0, 1).text = "数值"
    table.cell(0, 2).text = "风险等级"

    for idx, (_, row) in enumerate(indicators.iterrows(), start=1):
        table.cell(idx, 0).text = row["indicator_name"]
        table.cell(idx, 1).text = f"{row['indicator_value']:.2f}" if pd.notna(row["indicator_value"]) else "N/A"
        table.cell(idx, 2).text = row["risk_level"]
        fill = table.cell(idx, 2).fill
        color = RISK_RGB.get(row["risk_level"])
        if color:
            fill.solid()
            fill.fore_color.rgb = color


def export_ppt_report(
    metrics_df: pd.DataFrame,
    overall_df: pd.DataFrame,
    output_path: Path,
    assets_dir: Path,
    template_path: Path,
    ranking_df: pd.DataFrame,
) -> Path:
    assets_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(str(template_path)) if template_path.exists() else Presentation()

    for company in metrics_df["company_name"].unique():
        company_metrics = metrics_df[metrics_df["company_name"] == company]
        company_overall = overall_df[overall_df["company_name"] == company]
        score_map: dict[int, float] = dict(
            zip(company_overall["year"], company_overall["overall_risk_score"], strict=False)
        )

        slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])
        _add_title(slide1, f"{company} 公司概览")
        _add_textbox(
            slide1,
            f"总体风险分数: {score_map.get(company_metrics['year'].max(), float('nan')):.2f}",
            0.5,
            0.8,
            8.5,
            0.5,
        )
        _add_indicator_table(slide1, company_metrics)

        slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
        _add_title(slide2, f"{company} 趋势与排名")
        trend_path = assets_dir / f"{company}_trend.png"
        plot_trend(company_metrics, company, ["net_profit_margin", "current_ratio", "roe"], trend_path)
        slide2.shapes.add_picture(str(trend_path), Inches(0.5), Inches(1.5), width=Inches(5.5))

        latest_year = int(company_metrics["year"].max())
        bar_path = assets_dir / f"{company}_bar.png"
        plot_bar(metrics_df, "net_profit_margin", latest_year, bar_path)
        slide2.shapes.add_picture(str(bar_path), Inches(6.2), Inches(1.5), width=Inches(3.2))

        ranking_text = "\n".join(
            [
                f"{row['company_name']}: {row['indicator_value']:.2f}"
                for _, row in ranking_df.iterrows()
            ]
        )
        _add_textbox(slide2, f"{latest_year} 净利润率排名:\n{ranking_text}", 0.5, 4.7, 9, 1.5)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return output_path
